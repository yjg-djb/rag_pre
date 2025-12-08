from pathlib import Path
from typing import Tuple
import re
import os
import tempfile
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

# 安全导入 PyMuPDF（可选依赖）
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except Exception:
    fitz = None
    HAS_PYMUPDF = False


class DocumentDetector:
    """文档检测器 - 检测文档是否为纯文本"""
    
    def __init__(self):
        # 检测是否有 Office COM 支持（仅 Windows）
        self._has_office_com = False
        if os.name == 'nt':
            try:
                import win32com.client
                self._has_office_com = True
            except ImportError:
                pass
    
    def detect(self, file_path: str) -> Tuple[bool, str]:
        """
        检测文件是否为纯文本
        返回: (是否纯文本, 原因/描述)
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        
        if extension in ['.txt', '.md']:
            return self._detect_text_file(file_path)
        elif extension == '.docx':
            return self._detect_docx(file_path)
        elif extension == '.doc':
            return self._detect_doc(file_path)
        elif extension == '.xlsx':
            return self._detect_xlsx(file_path)
        elif extension == '.xls':
            return self._detect_xls(file_path)
        elif extension == '.pptx':
            return self._detect_pptx(file_path)
        elif extension == '.ppt':
            return self._detect_ppt(file_path)
        elif extension == '.pdf':
            return self._detect_pdf(file_path)
        else:
            return False, f"不支持的文件格式: {extension}"
    
    def _convert_old_format_to_new(self, file_path: str, target_format: str) -> str | None:
        """
        使用 COM 将旧格式转换为新格式
        
        Args:
            file_path: 源文件路径
            target_format: 目标格式 ('docx', 'xlsx', 'pptx')
        
        Returns:
            转换后的文件路径，失败返回None
        """
        if not self._has_office_com:
            return None
        
        # 安全导入 pywin32（可选依赖）
        try:
            import win32com.client
            import pythoncom
        except ImportError:
            # 内网未安装 pywin32 时，直接降级
            return None
        
        file_path = os.path.abspath(file_path)
        temp_dir = tempfile.gettempdir()
        base_name = Path(file_path).stem
        output_path = os.path.join(temp_dir, f"{base_name}_converted.{target_format}")
        
        try:
            pythoncom.CoInitialize()
            
            if target_format == 'docx':
                # Word 转换
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                word.DisplayAlerts = False
                
                doc = word.Documents.Open(file_path)
                doc.SaveAs(output_path, FileFormat=16)  # 16 = wdFormatXMLDocument (docx)
                doc.Close()
                word.Quit()
                
            elif target_format == 'xlsx':
                # Excel 转换
                excel = win32com.client.Dispatch("Excel.Application")
                excel.Visible = False
                excel.DisplayAlerts = False
                
                wb = excel.Workbooks.Open(file_path)
                wb.SaveAs(output_path, FileFormat=51)  # 51 = xlOpenXMLWorkbook (xlsx)
                wb.Close()
                excel.Quit()
                
            elif target_format == 'pptx':
                # PowerPoint 转换
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = 1
                powerpoint.DisplayAlerts = 0
                
                presentation = powerpoint.Presentations.Open(file_path, WithWindow=False)
                presentation.SaveAs(output_path, FileFormat=24)  # 24 = ppSaveAsOpenXMLPresentation (pptx)
                presentation.Close()
                powerpoint.Quit()
            
            pythoncom.CoUninitialize()
            
            if os.path.exists(output_path):
                return output_path
            return None
            
        except Exception as e:
            print(f"COM 转换失败: {e}")
            try:
                pythoncom.CoUninitialize()
            except:
                pass
            return None
    
    def _detect_doc(self, file_path: str) -> Tuple[bool, str]:
        """检测 .doc 文件 - 先转换为 docx 再检测"""
        # 尝试转换为 docx
        converted_path = self._convert_old_format_to_new(file_path, 'docx')
        
        if converted_path:
            try:
                result = self._detect_docx(converted_path)
                # 清理临时文件
                try:
                    os.unlink(converted_path)
                except:
                    pass
                return result
            except Exception as e:
                return False, f"处理 doc 文件错误: {str(e)}"
        else:
            return False, "不支持的格式: .doc (需要 Microsoft Word)"
    
    def _detect_xls(self, file_path: str) -> Tuple[bool, str]:
        """检测 .xls 文件 - 先转换为 xlsx 再检测"""
        # 尝试转换为 xlsx
        converted_path = self._convert_old_format_to_new(file_path, 'xlsx')
        
        if converted_path:
            try:
                result = self._detect_xlsx(converted_path)
                # 清理临时文件
                try:
                    os.unlink(converted_path)
                except:
                    pass
                return result
            except Exception as e:
                return False, f"处理 xls 文件错误: {str(e)}"
        else:
            return False, "不支持的格式: .xls (需要 Microsoft Excel)"
    
    def _detect_ppt(self, file_path: str) -> Tuple[bool, str]:
        """检测 .ppt 文件 - 先转换为 pptx 再检测"""
        # 尝试转换为 pptx
        converted_path = self._convert_old_format_to_new(file_path, 'pptx')
        
        if converted_path:
            try:
                result = self._detect_pptx(converted_path)
                # 清理临时文件
                try:
                    os.unlink(converted_path)
                except:
                    pass
                return result
            except Exception as e:
                return False, f"处理 ppt 文件错误: {str(e)}"
        else:
            return False, "不支持的格式: .ppt (需要 Microsoft PowerPoint)"
    
    def _detect_text_file(self, file_path: str) -> Tuple[bool, str]:
        """检测 txt/md 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 检查是否包含图片引用 ![...](...) 或 <img>
            if re.search(r'!\[.*?\]\(.*?\)', content) or '<img' in content.lower():
                return False, "包含图片引用"
            
            return True, "纯文本文件"
        except Exception as e:
            return False, f"读取文件错误: {str(e)}"
    
    def _detect_docx(self, file_path: str) -> Tuple[bool, str]:
        """检测 docx 文件 - 严格模式：只有纯文字才算纯文本"""
        doc = None
        try:
            # 打开 Word 文档
            doc = Document(file_path)
            
            # ========== 第1步：检查是否有表格 ==========
            # 只要有表格，就不是纯文本
            try:
                table_count = len(doc.tables)
                if table_count > 0:
                    return False, f"包含表格 ({table_count}个)"
            except Exception:
                # 访问表格失败，跳过
                pass
            
            # ========== 第2步：检查内嵌图片 ==========
            image_count = 0
            try:
                for rel in doc.part.rels.values():
                    try:
                        if "image" in rel.target_ref:
                            image_count += 1
                    except Exception:
                        continue
            except Exception:
                # 访问 relationships 失败
                pass
            
            if image_count > 0:
                return False, f"包含图片 ({image_count}个)"
            
            # ========== 第3步：检查嵌入对象（图表、形状等）==========
            embedded_count = 0
            try:
                for paragraph in doc.paragraphs:
                    try:
                        if hasattr(paragraph, '_element'):
                            for run in paragraph.runs:
                                try:
                                    if hasattr(run, '_element'):
                                        # 检查图片、图表等对象
                                        blips = run._element.xpath('.//a:blip')
                                        if blips:
                                            embedded_count += len(blips)
                                except Exception:
                                    continue
                    except Exception:
                        continue
            except Exception:
                # 遍历段落失败
                pass
            
            if embedded_count > 0:
                return False, f"包含嵌入对象 ({embedded_count}个)"
            
            # ========== 第4步：检查是否有文本内容 ==========
            has_text = False
            paragraph_count = 0
            try:
                for p in doc.paragraphs:
                    try:
                        if p.text and p.text.strip():
                            has_text = True
                            paragraph_count += 1
                    except Exception:
                        continue
            except Exception:
                pass
            
            if not has_text:
                return False, "文档无有效文本内容"
            
            # ========== 全部检查通过：纯文本文档 ==========
            return True, f"纯文本文档 ({paragraph_count}个段落)"
            
        except Exception as e:
            return False, f"处理 docx 错误: {str(e)}"
        finally:
            # 确保文档对象被释放
            if doc is not None:
                try:
                    # python-docx 的 Document 对象没有显式的 close 方法
                    # 但会在对象销毁时自动关闭
                    del doc
                except:
                    pass
    
    def _detect_xlsx(self, file_path: str) -> Tuple[bool, str]:
        """检测 xlsx 文件 - Excel本身就是表格格式，始终为富媒体"""
        wb = None
        try:
            wb = load_workbook(file_path)
            
            # Excel 文件本身就包含表格结构，不符合纯文本定义
            sheet_count = len(wb.worksheets)
            
            # 检查是否有图片、图表等
            has_images = False
            has_charts = False
            has_drawings = False
            
            try:
                for sheet in wb.worksheets:
                    try:
                        if hasattr(sheet, '_images') and len(sheet._images) > 0:
                            has_images = True
                        if hasattr(sheet, '_charts') and len(sheet._charts) > 0:
                            has_charts = True
                        if hasattr(sheet, '_drawings') and sheet._drawings:
                            has_drawings = True
                    except Exception:
                        continue
            except Exception:
                pass
            
            # 构建原因描述
            reasons = []
            if has_images:
                reasons.append("图片")
            if has_charts:
                reasons.append("图表")
            if has_drawings:
                reasons.append("绘图")
            
            if reasons:
                return False, f"Excel表格文件，包含{'+'.join(reasons)}"
            else:
                return False, f"Excel表格文件 ({sheet_count}个工作表)"
            
        except Exception as e:
            return False, f"处理 xlsx 错误: {str(e)}"
        finally:
            # 确保工作簿被关闭
            if wb is not None:
                try:
                    wb.close()
                except:
                    pass
    
    def _detect_pptx(self, file_path: str) -> Tuple[bool, str]:
        """检测 pptx 文件 - PPT本身就是富媒体格式，始终为富媒体"""
        prs = None
        try:
            # 打开 PPT 文件
            prs = Presentation(file_path)
            
            slide_count = len(prs.slides)
            
            # PPT 文件本身就包含幻灯片、布局等格式，不符合纯文本定义
            # 即使只有文字，也是带格式的演示文稿
            
            # 统计富媒体元素
            image_count = 0
            chart_count = 0
            shape_count = 0
            
            try:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        try:
                            if shape.shape_type == 13:  # PICTURE
                                image_count += 1
                            elif shape.shape_type == 3:  # CHART
                                chart_count += 1
                            elif hasattr(shape, 'image'):
                                shape_count += 1
                        except Exception:
                            # 某些 shape 可能无法访问，忽略
                            continue
            except Exception:
                # 如果遍历失败，仍然判定为富媒体
                pass
            
            # 构建原因描述
            reasons = []
            if image_count > 0:
                reasons.append(f"{image_count}张图片")
            if chart_count > 0:
                reasons.append(f"{chart_count}个图表")
            if shape_count > 0:
                reasons.append(f"{shape_count}个形状")
            
            if reasons:
                return False, f"PPT文件 ({slide_count}页)，包含{'+'.join(reasons)}"
            else:
                return False, f"PPT演示文稿 ({slide_count}页幻灯片)"
            
        except Exception as e:
            return False, f"处理 pptx 错误: {str(e)}"
        finally:
            # 确保文件被关闭（如果需要）
            if prs is not None:
                try:
                    # python-pptx 的 Presentation 对象没有显式的 close 方法
                    # 但会在对象销毁时自动关闭
                    del prs
                except:
                    pass
    
    def _detect_pdf(self, file_path: str) -> Tuple[bool, str]:
        """检测 pdf 文件 - 严格检测所有非文本元素"""
        if not HAS_PYMUPDF:
            return False, "PDF 检测未启用（未安装 PyMuPDF）"
        
        doc = None
        try:
            doc = fitz.open(file_path)
            page_count = len(doc)
            
            total_images = 0
            total_drawings = 0
            has_tables = False
            
            for page_num in range(page_count):
                page = doc[page_num]
                
                # ========== 检查图片 ==========
                image_list = page.get_images(full=True)
                total_images += len(image_list)
                
                # ========== 检查矢量图形和表格 ==========
                drawings = page.get_drawings()
                if drawings and len(drawings) > 0:
                    # PDF中的表格通常用线条绘制
                    # 如果有大量线条，很可能是表格
                    line_count = len([d for d in drawings if d.get('type') in ['l', 're']])
                    if line_count > 10:  # 超过10条线，可能是表格
                        has_tables = True
                    
                    # 复杂图形（非简单线条）
                    complex_drawings = [d for d in drawings if d.get('type') not in ['l', 're']]
                    total_drawings += len(complex_drawings)
            
            # ========== 判定逻辑 ==========
            if total_images > 0:
                return False, f"包含图片 ({total_images}个)"
            
            if has_tables:
                return False, "包含表格结构"
            
            if total_drawings > 0:
                return False, f"包含矢量图形 ({total_drawings}个)"
            
            # 全部检查通过：纯文本PDF
            return True, f"纯文本PDF ({page_count}页)"
            
        except Exception as e:
            return False, f"处理 PDF 错误: {str(e)}"
        finally:
            # 确保 PDF 文档被关闭
            if doc is not None:
                try:
                    doc.close()
                except:
                    pass
