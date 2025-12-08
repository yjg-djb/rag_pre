from pathlib import Path
import os
import tempfile
import subprocess
import shutil
import uuid
from typing import Optional, Dict
from datetime import datetime
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from utils.logger import get_logger

# 安全导入 PyMuPDF（可选依赖）
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

logger = get_logger("converter")


class DocumentConverter:
    """文档转换器 - 将各种格式转换为 docx"""
    
    def __init__(self, text_pipeline=None):
        """
        初始化转换器
        
        Args:
            text_pipeline: TextPipeline 实例（用于纯文本清洗与去重）
        """
        # 检测是否有 Office COM 支持（仅 Windows）
        self._has_office_com = False
        if os.name == 'nt':
            try:
                import win32com.client
                self._has_office_com = True
            except ImportError:
                pass
        
        # 检测 LibreOffice
        self.libreoffice_path = self._detect_libreoffice()
        if self.libreoffice_path:
            logger.info(f"LibreOffice 检测成功: {self.libreoffice_path}")
        else:
            logger.warning("LibreOffice 未检测到，将使用 Word COM 作为后备")
        
        # 文本管线（可选）
        self.text_pipeline = text_pipeline
    
    def _detect_libreoffice(self) -> Optional[str]:
        """检测 LibreOffice 安装路径"""
        try:
            from config import config
            
            # 1. 优先使用配置指定的路径
            if config.Conversion.LIBREOFFICE_PATH:
                if os.path.exists(config.Conversion.LIBREOFFICE_PATH):
                    return config.Conversion.LIBREOFFICE_PATH
            
            # 2. 尝试常见安装路径
            for path in config.Conversion.LIBREOFFICE_DEFAULT_PATHS:
                if os.path.exists(path):
                    return path
            
            # 3. 尝试从 PATH 环境变量查找
            soffice = shutil.which("soffice")
            if soffice:
                return soffice
            
            return None
        except:
            return None
    
    def convert_to_docx(self, input_file: str, output_file: str, doc_name: str = "unknown", apply_pipeline: bool = True) -> Dict:
        """
        将文件转换为目标格式（集成文本管线）
        
        Args:
            input_file: 输入文件路径
            output_file: 输出文件路径
            doc_name: 文档名称（用于日志）
            apply_pipeline: 是否应用文本管线（仅纯文本才应用）
        
        Returns:
            {
                "success": bool,
                "message": str,
                "pipeline_stats": dict (如果使用了文本管线)
            }
        """
        path = Path(input_file)
        extension = path.suffix.lower()
        output_ext = Path(output_file).suffix.lower()
        
        try:
            # 先转换为中间格式或提取文本
            if extension in ['.txt', '.md']:
                success = self._txt_to_docx(input_file, output_file)
            elif extension == '.docx':
                # docx 转 docx（复制后检测是否需要清洗）
                if output_ext == '.docx':
                    success = self._copy_docx(input_file, output_file)
                else:
                    success = False
            elif extension == '.doc':
                # .doc 转 .docx
                success = self._doc_to_docx(input_file, output_file)
            elif extension == '.xlsx':
                # xlsx 转为 docx（用于纯文本检测）或保持 xlsx
                if output_ext == '.docx':
                    success = self._xlsx_to_docx(input_file, output_file)
                elif output_ext == '.xlsx':
                    success = self._copy_file(input_file, output_file)
                else:
                    success = False
            elif extension == '.xls':
                # .xls 转 .xlsx 或 .docx
                if output_ext == '.xlsx':
                    success = self._xls_to_xlsx(input_file, output_file)
                elif output_ext == '.docx':
                    # 先转xlsx再转docx
                    temp_xlsx = self._convert_old_to_new(input_file, 'xlsx')
                    if temp_xlsx:
                        success = self._xlsx_to_docx(temp_xlsx, output_file)
                        try:
                            os.remove(temp_xlsx)
                        except:
                            pass
                    else:
                        success = False
                else:
                    success = False
            elif extension == '.pptx':
                # pptx 转为 docx（用于纯文本检测）或保持 pptx
                if output_ext == '.docx':
                    success = self._pptx_to_docx(input_file, output_file)
                elif output_ext == '.pptx':
                    success = self._copy_file(input_file, output_file)
                else:
                    success = False
            elif extension == '.ppt':
                # .ppt 转 .pptx 或 .docx
                if output_ext == '.pptx':
                    success = self._ppt_to_pptx(input_file, output_file)
                elif output_ext == '.docx':
                    # 先转pptx再转docx
                    temp_pptx = self._convert_old_to_new(input_file, 'pptx')
                    if temp_pptx:
                        success = self._pptx_to_docx(temp_pptx, output_file)
                        try:
                            os.remove(temp_pptx)
                        except:
                            pass
                    else:
                        success = False
                else:
                    success = False
            elif extension == '.pdf':
                success = self._pdf_to_docx(input_file, output_file)
            else:
                return {"success": False, "message": f"不支持的格式: {extension}"}
            
            if not success:
                return {"success": False, "message": "格式转换失败"}
            
            # 如果有文本管线且输出是 docx，且明确要求应用管线（仅纯文本）
            if self.text_pipeline and apply_pipeline and os.path.exists(output_file) and output_ext == '.docx':
                logger.info(f"[{doc_name}] 应用文本管线进行清洗与去重")
                
                # 从生成的 docx 提取文本
                text = self._extract_text_from_docx(output_file)
                
                # 应用文本管线
                result = self.text_pipeline.process(text, doc_name)
                
                # 无论是否去重，都要保存清洗后的文本
                if result.get("cleaned_text"):
                    cleaned_text = result["cleaned_text"]
                    self._write_cleaned_text_to_docx(cleaned_text, output_file)
                    logger.debug(f"[{doc_name}] 清洗后的文本已保存到: {output_file}")
                
                if not result["success"]:
                    # 文档级去重命中（但清洗后的文件已保存）
                    return {
                        "success": False,
                        "message": result["message"],
                        "pipeline_stats": result["stats"],
                        "doc_duplicate": result.get("doc_duplicate", False)
                    }
                
                logger.info(f"[{doc_name}] 文本管线处理完成: {result['stats']}")
                
                return {
                    "success": True,
                    "message": "转换与清洗成功",
                    "pipeline_stats": result["stats"]
                }
            
            return {"success": True, "message": "转换成功"}
            
        except Exception as e:
            logger.error(f"转换错误: {e}", exc_info=True)
            return {"success": False, "message": f"转换错误: {str(e)}"}
    
    def _convert_old_to_new(self, input_file: str, target_format: str) -> str | None:
        """
        使用 LibreOffice 或 COM 将旧格式转换为新格式
        
        Args:
            input_file: 源文件路径
            target_format: 目标格式 ('docx', 'xlsx', 'pptx')
        
        Returns:
            转换后的文件路径，失败返回None
        """
        try:
            from config import config
            backend = config.Conversion.BACKEND
        except:
            backend = "auto"
        
        # 优先级策略
        if backend == "libreoffice":
            # 只使用 LibreOffice
            result = self._convert_with_libreoffice(input_file, target_format)
            if result:
                return result
        elif backend == "word":
            # 只使用 Word COM
            result = self._convert_with_word_com(input_file, target_format)
            if result:
                return result
        else:  # auto
            # 优先 LibreOffice，失败时尝试 Word
            result = self._convert_with_libreoffice(input_file, target_format)
            if result:
                return result
            logger.debug(f"LibreOffice 转换失败，尝试 Word COM")
            result = self._convert_with_word_com(input_file, target_format)
            if result:
                return result
        
        return None
    
    def _convert_with_libreoffice(self, input_file: str, target_format: str) -> str | None:
        """
        使用 LibreOffice 转换文档
        
        Args:
            input_file: 源文件路径
            target_format: 目标格式 ('docx', 'xlsx', 'pptx')
        
        Returns:
            转换后的文件路径，失败返回None
        
        修复:
        1. 使用 UUID 生成临时文件名，避免中文路径问题
        2. 使用项目内临时目录，避免系统临时目录权限问题
        """
        if not self.libreoffice_path:
            return None
        
        input_file = os.path.abspath(input_file)
        
        # 修复：使用项目内的临时目录，避免系统临时目录权限问题
        project_temp_dir = Path("storage/temp")
        project_temp_dir.mkdir(parents=True, exist_ok=True)
        temp_dir = str(project_temp_dir.absolute())
        
        # 修复：使用 UUID 作为临时文件名，避免中文路径问题
        import uuid
        temp_filename = f"temp_{uuid.uuid4().hex[:12]}.{target_format}"
        expected_output = os.path.join(temp_dir, temp_filename)
        
        # 映射格式到 LibreOffice 过滤器
        format_map = {
            'docx': 'MS Word 2007 XML',
            'xlsx': 'Calc MS Excel 2007 XML',
            'pptx': 'Impress MS PowerPoint 2007 XML'
        }
        
        filter_name = format_map.get(target_format)
        if not filter_name:
            return None
        
        # 获取超时配置（移至 try 外）
        timeout = 60  # 默认值
        try:
            from config import config
            timeout = config.Conversion.CONVERSION_TIMEOUT
        except:
            pass
        
        try:
            
            # 构建命令
            cmd = [
                self.libreoffice_path,
                '--headless',
                '--nologo',
                '--nolockcheck',
                '--convert-to', f'{target_format}:"{filter_name}"',
                '--outdir', temp_dir,
                input_file
            ]
            
            logger.debug(f"LibreOffice 转换命令: {' '.join(cmd)}")
            logger.debug(f"预期输出文件: {expected_output}")
            
            # 执行转换
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=timeout,
                creationflags=subprocess.CREATE_NO_WINDOW if os.name == 'nt' else 0
            )
            
            # 修复：LibreOffice 实际输出文件名可能不同（使用原文件名）
            # 需要查找实际生成的文件
            base_name = Path(input_file).stem
            actual_output = os.path.join(temp_dir, f"{base_name}.{target_format}")
            
            # 检查两种可能的输出文件
            if os.path.exists(expected_output):
                logger.debug(f"LibreOffice 转换成功（UUID文件名）: {expected_output}")
                return expected_output
            elif os.path.exists(actual_output):
                # LibreOffice 使用了原始文件名，重命名为 UUID 文件名
                logger.debug(f"LibreOffice 使用了原始文件名: {actual_output}，重命名为: {expected_output}")
                try:
                    shutil.move(actual_output, expected_output)
                    return expected_output
                except Exception as e:
                    logger.warning(f"重命名文件失败: {e}，返回原文件")
                    return actual_output
            elif result.returncode == 0:
                # 转换成功但文件名不匹配，搜索临时目录
                logger.warning(f"LibreOffice 转换成功但找不到预期文件，搜索临时目录")
                for file in Path(temp_dir).glob(f"*.{target_format}"):
                    if file.stat().st_mtime > (datetime.now().timestamp() - 10):  # 10秒内创建的文件
                        logger.info(f"找到最近创建的文件: {file}")
                        try:
                            shutil.move(str(file), expected_output)
                            return expected_output
                        except:
                            return str(file)
                logger.error(f"LibreOffice 转换成功但找不到输出文件")
                return None
            else:
                logger.warning(f"LibreOffice 转换失败: {result.stderr}")
                return None
                
        except subprocess.TimeoutExpired:
            logger.error(f"LibreOffice 转换超时（{timeout}秒）")
            return None
        except Exception as e:
            logger.debug(f"LibreOffice 转换异常: {e}")
            return None
    
    def _convert_with_word_com(self, input_file: str, target_format: str) -> str | None:
        """
        使用 Word COM 转换文档（仅支持 docx）
        
        Args:
            input_file: 源文件路径
            target_format: 目标格式 ('docx', 'xlsx', 'pptx')
        
        Returns:
            转换后的文件路径，失败返回None
        """
        if not self._has_office_com:
            return None
        
        # 安全导入 pywin32（可选依赖）
        try:
            import win32com.client
            import pythoncom
        except ImportError:
            # 内网未安装 pywin32 时，直接降级
            return None
        
        input_file = os.path.abspath(input_file)
        temp_dir = tempfile.gettempdir()
        base_name = Path(input_file).stem
        output_path = os.path.join(temp_dir, f"{base_name}_temp.{target_format}")
        
        try:
            pythoncom.CoInitialize()
            
            if target_format == 'docx':
                # Word 转换
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                word.DisplayAlerts = False
                
                doc = word.Documents.Open(input_file)
                doc.SaveAs(output_path, FileFormat=16)  # 16 = wdFormatXMLDocument
                doc.Close()
                word.Quit()
                
            elif target_format == 'xlsx':
                # Excel 转换
                excel = win32com.client.Dispatch("Excel.Application")
                excel.Visible = False
                excel.DisplayAlerts = False
                
                wb = excel.Workbooks.Open(input_file)
                wb.SaveAs(output_path, FileFormat=51)  # 51 = xlOpenXMLWorkbook
                wb.Close()
                excel.Quit()
                
            elif target_format == 'pptx':
                # PowerPoint 转换
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = 1
                powerpoint.DisplayAlerts = 0
                
                presentation = powerpoint.Presentations.Open(input_file, WithWindow=False)
                presentation.SaveAs(output_path, FileFormat=24)  # 24 = ppSaveAsOpenXMLPresentation
                presentation.Close()
                powerpoint.Quit()
            
            pythoncom.CoUninitialize()
            
            if os.path.exists(output_path):
                return output_path
            return None
            
        except Exception as e:
            print(f"COM 转换失败: {e}")
            try:
                pythoncom.CoUninitialize()
            except:
                pass
            return None
    
    def _doc_to_docx(self, input_file: str, output_file: str) -> bool:
        """将 .doc 转换为 .docx"""
        # 优先使用 LibreOffice/Word COM 的通用转换方法
        temp_docx = self._convert_old_to_new(input_file, 'docx')
        if temp_docx:
            try:
                import shutil
                shutil.move(temp_docx, output_file)
                logger.debug(f"DOC 转 DOCX 成功: {output_file}")
                return True
            except Exception as e:
                logger.error(f"DOC 转 DOCX 移动文件错误: {e}")
                return False
        else:
            logger.error("DOC 转 DOCX 失败：未检测到 LibreOffice 或 Word")
            return False
    
    def _xls_to_docx(self, input_file: str, output_file: str) -> bool:
        """将 .xls 转换为 .docx"""
        # 先转为 xlsx，再转为 docx
        temp_xlsx = self._convert_old_to_new(input_file, 'xlsx')
        if temp_xlsx:
            try:
                result = self._xlsx_to_docx(temp_xlsx, output_file)
                # 清理临时文件
                try:
                    os.unlink(temp_xlsx)
                except:
                    pass
                return result
            except Exception as e:
                print(f"XLS 转换错误: {e}")
                return False
        else:
            print("需要 Microsoft Excel 支持 .xls 格式")
            return False
    
    def _ppt_to_docx(self, input_file: str, output_file: str) -> bool:
        """将 .ppt 转换为 .docx"""
        # 先转为 pptx，再转为 docx
        temp_pptx = self._convert_old_to_new(input_file, 'pptx')
        if temp_pptx:
            try:
                result = self._pptx_to_docx(temp_pptx, output_file)
                # 清理临时文件
                try:
                    os.unlink(temp_pptx)
                except:
                    pass
                return result
            except Exception as e:
                print(f"PPT 转换错误: {e}")
                return False
        else:
            print("需要 Microsoft PowerPoint 支持 .ppt 格式")
            return False
    
    def _txt_to_docx(self, input_file: str, output_file: str) -> bool:
        """TXT/MD 转 DOCX"""
        try:
            with open(input_file, 'r', encoding='utf-8') as f:
                content = f.read()
            
            doc = Document()
            
            # 按段落分割
            paragraphs = content.split('\n\n')
            for paragraph in paragraphs:
                if paragraph.strip():
                    # 检查是否是标题（以 # 开头的 Markdown）
                    if paragraph.strip().startswith('#'):
                        level = len(paragraph) - len(paragraph.lstrip('#'))
                        text = paragraph.lstrip('#').strip()
                        doc.add_heading(text, level=min(level, 3))
                    else:
                        # 处理单行换行
                        lines = paragraph.split('\n')
                        para = doc.add_paragraph()
                        for i, line in enumerate(lines):
                            if i > 0:
                                para.add_run('\n')
                            para.add_run(line)
            
            doc.save(output_file)
            return True
        except Exception as e:
            print(f"TXT 转换错误: {e}")
            return False
    
    def _copy_docx(self, input_file: str, output_file: str) -> bool:
        """复制 DOCX 文件"""
        try:
            import shutil
            shutil.copy2(input_file, output_file)
            return True
        except Exception as e:
            print(f"DOCX 复制错误: {e}")
            return False
    
    def _copy_file(self, input_file: str, output_file: str) -> bool:
        """通用文件复制"""
        try:
            import shutil
            shutil.copy2(input_file, output_file)
            return True
        except Exception as e:
            print(f"文件复制错误: {e}")
            return False
    
    def _xls_to_xlsx(self, input_file: str, output_file: str) -> bool:
        """将 .xls 转换为 .xlsx"""
        temp_xlsx = self._convert_old_to_new(input_file, 'xlsx')
        if temp_xlsx:
            try:
                import shutil
                shutil.move(temp_xlsx, output_file)
                return True
            except Exception as e:
                print(f"XLS 转 XLSX 错误: {e}")
                return False
        else:
            logger.error("XLS 转 XLSX 失败：未检测到 LibreOffice 或 Excel")
            return False
    
    def _ppt_to_pptx(self, input_file: str, output_file: str) -> bool:
        """将 .ppt 转换为 .pptx"""
        temp_pptx = self._convert_old_to_new(input_file, 'pptx')
        if temp_pptx:
            try:
                import shutil
                shutil.move(temp_pptx, output_file)
                return True
            except Exception as e:
                print(f"PPT 转 PPTX 错误: {e}")
                return False
        else:
            logger.error("PPT 转 PPTX 失败：未检测到 LibreOffice 或 PowerPoint")
            return False
    
    def _xlsx_to_docx(self, input_file: str, output_file: str) -> bool:
        """XLSX 转 DOCX - 增强异常处理"""
        wb = None
        try:
            wb = load_workbook(input_file, data_only=True)
            doc = Document()
            
            for sheet in wb.worksheets:
                # 添加工作表标题
                doc.add_heading(sheet.title, level=1)
                
                # 获取所有行
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    continue
                
                # 创建表格
                max_cols = max(len(row) for row in rows)
                table = doc.add_table(rows=len(rows), cols=max_cols)
                table.style = 'Light Grid Accent 1'
                
                # 填充数据
                for i, row in enumerate(rows):
                    for j, cell_value in enumerate(row):
                        if cell_value is not None:
                            table.rows[i].cells[j].text = str(cell_value)
                
                doc.add_paragraph()  # 添加空行
            
            doc.save(output_file)
            return True
            
        except Exception as e:
            print(f"XLSX 转换错误: {e}")
            return False
        finally:
            # 确保工作簿被关闭
            if wb is not None:
                try:
                    wb.close()
                except:
                    pass
    
    def _pptx_to_docx(self, input_file: str, output_file: str) -> bool:
        """PPTX 转 DOCX - 增强异常处理"""
        prs = None
        try:
            # 打开 PPT 文件
            prs = Presentation(input_file)
            doc = Document()
            
            for i, slide in enumerate(prs.slides):
                # 添加幻灯片标题
                doc.add_heading(f'Slide {i + 1}', level=1)
                
                # 提取文本
                try:
                    for shape in slide.shapes:
                        try:
                            # 尝试提取文本
                            if hasattr(shape, 'text'):
                                text = shape.text
                                if text and text.strip():
                                    # 判断是否是标题
                                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                                        doc.add_heading(text, level=2)
                                    else:
                                        doc.add_paragraph(text)
                        except Exception:
                            # 单个 shape 处理失败，继续下一个
                            continue
                        
                        # 处理表格
                        try:
                            if hasattr(shape, 'has_table') and shape.has_table:
                                table = shape.table
                                doc_table = doc.add_table(rows=len(table.rows), cols=len(table.columns))
                                doc_table.style = 'Light Grid Accent 1'
                                
                                for row_idx, row in enumerate(table.rows):
                                    for col_idx, cell in enumerate(row.cells):
                                        doc_table.rows[row_idx].cells[col_idx].text = cell.text
                        except Exception:
                            # 表格处理失败，继续
                            continue
                except Exception:
                    # 整个幻灯片处理失败，继续下一张
                    pass
                
                doc.add_paragraph()  # 幻灯片间添加空行
            
            doc.save(output_file)
            return True
            
        except Exception as e:
            print(f"PPTX 转换错误: {e}")
            return False
        finally:
            # 确保文件被关闭
            if prs is not None:
                try:
                    del prs
                except:
                    pass
    
    def _pdf_to_docx(self, input_file: str, output_file: str) -> bool:
        """PDF 转 DOCX - 增强异常处理"""
        if fitz is None:
            print("未安装 PyMuPDF，无法将 PDF 转为 DOCX")
            return False
        
        pdf_doc = None
        try:
            pdf_doc = fitz.open(input_file)
            doc = Document()
            
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                
                # 添加页面标题
                if page_num > 0:
                    doc.add_page_break()
                doc.add_heading(f'Page {page_num + 1}', level=1)
                
                # 提取文本
                text = page.get_text()
                
                # 按段落分割
                paragraphs = text.split('\n\n')
                for paragraph in paragraphs:
                    if paragraph.strip():
                        doc.add_paragraph(paragraph.strip())
            
            doc.save(output_file)
            return True
            
        except Exception as e:
            print(f"PDF 转换错误: {e}")
            return False
        finally:
            # 确保 PDF 文档被关闭
            if pdf_doc is not None:
                try:
                    pdf_doc.close()
                except:
                    pass
    
    def _extract_text_from_docx(self, docx_file: str) -> str:
        """
        从 DOCX 文件提取纯文本
        
        Args:
            docx_file: DOCX 文件路径
        
        Returns:
            提取的文本
        """
        try:
            doc = Document(docx_file)
            paragraphs = []
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            
            return '\n\n'.join(paragraphs)
            
        except Exception as e:
            logger.error(f"提取 DOCX 文本失败: {e}")
            return ""
    
    def _write_cleaned_text_to_docx(self, text: str, output_file: str) -> bool:
        """
        将清洗后的文本写入 DOCX 文件
        
        Args:
            text: 清洗后的文本
            output_file: 输出文件路径
        
        Returns:
            是否成功
        """
        try:
            doc = Document()
            
            # 按双换行拆分段落
            paragraphs = text.split('\n\n')
            
            for para_text in paragraphs:
                if para_text.strip():
                    # 检查是否是 Markdown 标题
                    if para_text.strip().startswith('#'):
                        level = len(para_text) - len(para_text.lstrip('#'))
                        heading_text = para_text.lstrip('#').strip()
                        doc.add_heading(heading_text, level=min(level, 3))
                    else:
                        # 处理单行换行（保留段内换行）
                        lines = para_text.split('\n')
                        para = doc.add_paragraph()
                        for i, line in enumerate(lines):
                            if i > 0:
                                para.add_run('\n')
                            para.add_run(line)
            
            doc.save(output_file)
            return True
            
        except Exception as e:
            logger.error(f"写入 DOCX 文件失败: {e}")
            return False
